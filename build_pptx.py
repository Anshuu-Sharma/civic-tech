#!/usr/bin/env python3
"""
Build the JanSunwai AI hackathon PPTX with content + rendered Mermaid diagrams.
Uses python-pptx to build the presentation.
Diagrams are pre-rendered PNGs in the diagrams/ directory.
"""

import os
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── Configuration ───
OUTPUT_FILE = "JanSunwai_AI_Hackathon_Submission.pptx"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
DIAGRAM_DIR = Path("diagrams")

# Color palette
BG_DARK    = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD    = RGBColor(0x1A, 0x24, 0x3B)
BG_CARD2   = RGBColor(0x15, 0x1F, 0x34)
BLUE       = RGBColor(0x38, 0xBD, 0xF8)
GREEN      = RGBColor(0x4A, 0xDE, 0x80)
ORANGE     = RGBColor(0xFB, 0xBF, 0x24)
RED        = RGBColor(0xF8, 0x71, 0x71)
PURPLE     = RGBColor(0xA7, 0x8B, 0xFA)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY      = RGBColor(0xCB, 0xD5, 0xE1)
MGRAY      = RGBColor(0x94, 0xA3, 0xB8)
DBLUE      = RGBColor(0x1E, 0x40, 0xAF)


# ─── Helpers ───

def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def card(slide, l, t, w, h, color=BG_CARD):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return s

def rect(slide, l, t, w, h, color=BG_CARD):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return s

def txt(slide, l, t, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, name="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold
    p.font.name = name; p.alignment = align
    return tb

def mtxt(slide, l, t, w, h, lines, name="Calibri"):
    """Multi-line text box. lines = list of dicts with text, size, color, bold, align, space_after."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, ld in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ld.get("text", "")
        p.font.size = Pt(ld.get("size", 16))
        p.font.color.rgb = ld.get("color", LGRAY)
        p.font.bold = ld.get("bold", False)
        p.font.name = name
        p.alignment = ld.get("align", PP_ALIGN.LEFT)
        if "space_after" in ld: p.space_after = Pt(ld["space_after"])
        if "space_before" in ld: p.space_before = Pt(ld["space_before"])
    return tb

def accent_line(slide, l, t, w, color=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(3))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def dot(slide, l, t, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, Inches(0.18), Inches(0.18))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()

def add_img(slide, name, left, top, max_w, max_h):
    """Add diagram PNG maintaining aspect ratio within max_w x max_h bounds."""
    path = DIAGRAM_DIR / f"{name}.png"
    if not path.exists():
        card(slide, left, top, max_w, max_h, RGBColor(0x1E, 0x29, 0x3B))
        txt(slide, left, top + max_h // 2, max_w, Inches(0.4), f"[{name}]", 14, MGRAY, align=PP_ALIGN.CENTER)
        return

    img = Image.open(path)
    iw, ih = img.size
    aspect = iw / ih

    # Fit within bounds keeping aspect ratio
    target_w = max_w
    target_h = int(target_w / aspect) if aspect > 0 else max_h

    if target_h > max_h:
        target_h = max_h
        target_w = int(max_h * aspect)

    # Center within the bounding box
    offset_x = left + (max_w - target_w) // 2
    offset_y = top + (max_h - target_h) // 2

    slide.shapes.add_picture(str(path), offset_x, offset_y, target_w, target_h)


def slide_header(slide, title, subtitle=None):
    """Standard slide header with title + accent line."""
    txt(slide, Inches(0.8), Inches(0.35), Inches(10), Inches(0.6), title, 34, BLUE, bold=True)
    accent_line(slide, Inches(0.8), Inches(0.95), Inches(3.5), BLUE)
    if subtitle:
        txt(slide, Inches(0.8), Inches(1.05), Inches(10), Inches(0.4), subtitle, 14, MGRAY)


# ═══════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════

def slide_01_title(prs):
    """Slide 1: Title"""
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)

    # Top accent
    accent_line(sl, Inches(0), Inches(0), SLIDE_W, BLUE)

    # Title
    txt(sl, Inches(1), Inches(1.3), Inches(11.3), Inches(1),
        "JanSunwai AI", 56, BLUE, bold=True, align=PP_ALIGN.CENTER)

    txt(sl, Inches(1.5), Inches(2.5), Inches(10.3), Inches(0.7),
        "AI-Powered Civic Grievance Resolution Platform", 28, WHITE, align=PP_ALIGN.CENTER)

    txt(sl, Inches(2), Inches(3.3), Inches(9.3), Inches(0.5),
        '"Your Voice Matters. We Make It Heard."', 20, GREEN, align=PP_ALIGN.CENTER)

    accent_line(sl, Inches(5.5), Inches(4.1), Inches(2.3), ORANGE)

    # Info cards
    cy = Inches(4.6)
    infos = [
        (Inches(1.5), "TEAM NAME", BLUE, "JanSunwai AI"),
        (Inches(5.1), "PROBLEM STATEMENT", ORANGE, "AI Civic Grievance Resolution"),
        (Inches(8.9), "TEAM LEADER", GREEN, "[Your Name]"),
    ]
    for cx, label, clr, value in infos:
        card(sl, cx, cy, Inches(3.1), Inches(1.3))
        txt(sl, cx + Inches(0.15), cy + Inches(0.12), Inches(2.8), Inches(0.3),
            label, 11, clr, bold=True, align=PP_ALIGN.CENTER)
        txt(sl, cx + Inches(0.15), cy + Inches(0.5), Inches(2.8), Inches(0.6),
            value, 20, WHITE, bold=True, align=PP_ALIGN.CENTER)

    txt(sl, Inches(3), Inches(6.5), Inches(7.3), Inches(0.5),
        "AWS AI for Bharat Hackathon 2025", 14, MGRAY, align=PP_ALIGN.CENTER)


def slide_02_idea(prs):
    """Slide 2: Brief About the Idea"""
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    slide_header(sl, "Brief About the Idea")

    # Left: description
    mtxt(sl, Inches(0.8), Inches(1.3), Inches(7.2), Inches(3), [
        {"text": "JanSunwai AI is an end-to-end, AI-powered civic grievance resolution platform for Indian citizens.",
         "size": 18, "color": WHITE, "bold": True, "space_after": 10},
        {"text": "It transforms the broken complaint redressal system by combining multimodal complaint filing "
                 "(web form + Hindi voice agent), AI-driven classification & routing, automated 5-level escalation, "
                 "and public transparency dashboards.",
         "size": 14, "color": LGRAY, "space_after": 10},
        {"text": "A citizen in rural India can call and speak in Hindi to file a complaint — no forms, no literacy barrier. "
                 "AI automatically classifies the issue, scores severity, detects duplicates, routes it to the right department, "
                 "and generates legal rights awareness for the citizen.",
         "size": 14, "color": LGRAY, "space_after": 8},
        {"text": "Built for scale: 10+ Indian languages, PostGIS ward-level routing, and an hourly auto-escalation cron "
                 "that ensures every single complaint reaches resolution.",
         "size": 14, "color": LGRAY},
    ])

    # Right: stats cards
    stats = [
        ("2.8 Cr", "Pending civic\ngrievances across India", RED),
        ("45 Days", "Average resolution\ntime for complaints", ORANGE),
        ("67%", "Citizens unaware of\ntheir legal rights", BLUE),
        ("12+", "Departments with\nzero communication", GREEN),
    ]
    for i, (num, label, clr) in enumerate(stats):
        cy = Inches(0.45) + i * Inches(1.7)
        card(sl, Inches(8.4), cy, Inches(4.4), Inches(1.45))
        txt(sl, Inches(8.7), cy + Inches(0.15), Inches(2), Inches(0.7), num, 34, clr, bold=True)
        txt(sl, Inches(10.6), cy + Inches(0.3), Inches(2), Inches(0.8), label, 13, LGRAY)

    # Bottom: capability cards
    caps = [
        ("Voice-First", "Hindi/English voice\ncomplaint filing", BLUE),
        ("AI Classification", "Gemini-powered\ncategory + severity", ORANGE),
        ("Auto-Escalation", "5-level SLA-enforced\nescalation ladder", RED),
        ("Public Dashboards", "Ward-level heatmaps\n& scorecards", GREEN),
        ("Legal Awareness", "AI-generated rights\n& RTI support", PURPLE),
    ]
    for i, (title, desc, clr) in enumerate(caps):
        cx = Inches(0.4) + i * Inches(2.55)
        card(sl, cx, Inches(5.1), Inches(2.4), Inches(2))
        txt(sl, cx + Inches(0.15), Inches(5.25), Inches(2.1), Inches(0.35),
            title, 14, clr, bold=True, align=PP_ALIGN.CENTER)
        txt(sl, cx + Inches(0.15), Inches(5.7), Inches(2.1), Inches(1),
            desc, 12, LGRAY, align=PP_ALIGN.CENTER)


def slide_03_usp(prs):
    """Slide 3: Differentiation & USP"""
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    slide_header(sl, "How is it Different? How Does it Solve the Problem?")

    # Comparison table
    headers = ["Aspect", "Existing Systems", "JanSunwai AI"]
    hcolors = [DBLUE, DBLUE, RGBColor(0x05, 0x96, 0x69)]
    cx_list = [Inches(0.6), Inches(3.2), Inches(7.0)]
    cw_list = [Inches(2.5), Inches(3.7), Inches(5.5)]
    hy = Inches(1.3)

    for hdr, hc, cx, cw in zip(headers, hcolors, cx_list, cw_list):
        rect(sl, cx, hy, cw, Inches(0.45), hc)
        txt(sl, cx + Inches(0.1), hy + Inches(0.05), cw - Inches(0.2), Inches(0.35),
            hdr, 13, WHITE, bold=True, align=PP_ALIGN.CENTER)

    rows = [
        ("Filing Method",    "Web forms only",            "Voice Agent (Hindi) + Web Form + 10 Languages"),
        ("Classification",   "Manual / keyword-based",    "AI-powered (Gemini 2.0 Flash + Vision)"),
        ("Routing",          "Manual department assignment","Auto-route to dept + ward officer (PostGIS)"),
        ("Escalation",       "None / manual follow-up",   "5-level auto-escalation with SLA tracking"),
        ("Transparency",     "Opaque — no public data",   "Public heatmaps + ward scorecards"),
        ("Legal Rights",     "None",                      "AI-generated legal summaries per category"),
        ("Duplicates",       "None",                      "Semantic + geographic community clustering"),
        ("Accessibility",    "Requires literacy + internet","Voice-first — works for low-literacy citizens"),
    ]

    for i, (asp, old, new) in enumerate(rows):
        ry = hy + Inches(0.5) + i * Inches(0.52)
        bg = BG_CARD if i % 2 == 0 else BG_CARD2
        for cx, cw, text, clr in zip(cx_list, cw_list,
                                      [asp, old, new], [ORANGE, MGRAY, GREEN]):
            rect(sl, cx, ry, cw, Inches(0.48), bg)
            txt(sl, cx + Inches(0.1), ry + Inches(0.07), cw - Inches(0.2), Inches(0.34),
                text, 12, clr, bold=(clr == ORANGE))

    # USP box
    uy = Inches(5.7)
    card(sl, Inches(0.6), uy, Inches(11.9), Inches(1.5), DBLUE)
    txt(sl, Inches(0.9), uy + Inches(0.1), Inches(2), Inches(0.3), "USP", 16, ORANGE, bold=True)
    txt(sl, Inches(0.9), uy + Inches(0.45), Inches(11.3), Inches(0.9),
        "The only platform that combines AI voice filing in Indian languages, automated severity-based "
        "escalation, and public ward-level accountability dashboards — making civic grievance resolution "
        "accessible, intelligent, and transparent.", 15, WHITE)


def slide_04_features(prs):
    """Slide 4: Features + Mindmap"""
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    slide_header(sl, "Key Features")

    features = [
        ("Multi-Modal Filing",     "Web wizard + Hindi voice agent + photo upload",   BLUE),
        ("AI Intelligence Engine", "Gemini classification, vision, severity scoring",  ORANGE),
        ("Smart Routing",          "Auto-assign to dept + ward officer via PostGIS",   GREEN),
        ("5-Level Escalation",     "Auto-escalate to Commissioner if SLA breached",    RED),
        ("Public Dashboards",      "Ward heatmaps, category stats, scorecards",        BLUE),
        ("Admin Portal",           "Role-based queue, timelines, manual controls",     ORANGE),
        ("Legal Awareness",        "AI-generated rights summaries + RTI drafting",     GREEN),
        ("Citizen Verification",   "SMS/WhatsApp verify + satisfaction scoring",       RED),
    ]

    for i, (title, desc, clr) in enumerate(features):
        fy = Inches(1.15) + i * Inches(0.76)
        dot(sl, Inches(0.8), fy + Inches(0.08), clr)
        txt(sl, Inches(1.15), fy - Inches(0.03), Inches(4.8), Inches(0.3), title, 15, WHITE, bold=True)
        txt(sl, Inches(1.15), fy + Inches(0.26), Inches(4.8), Inches(0.3), desc, 11, MGRAY)

    # Feature mindmap (right half)
    add_img(sl, "feature_mindmap", Inches(6), Inches(1.0), Inches(7), Inches(6.2))


def slide_05_process_flow(prs):
    """Slide 5: Process Flow Diagram (full slide)"""
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    slide_header(sl, "Process Flow — Grievance Lifecycle")

    # The process flow is very tall, so show it scaled to fit
    add_img(sl, "process_flow", Inches(3), Inches(1.2), Inches(7), Inches(6))

    # Add summary steps on the left
    steps = [
        ("1", "Citizen files via Web or Voice", BLUE),
        ("2", "AI classifies + scores severity", ORANGE),
        ("3", "PostGIS resolves ward", GREEN),
        ("4", "Duplicate detection + clustering", PURPLE),
        ("5", "Auto-route to department", BLUE),
        ("6", "Generate complaint # + legal rights", GREEN),
        ("7", "Officer acts within SLA", ORANGE),
        ("8", "Auto-escalate if breached", RED),
        ("9", "Citizen verifies resolution", GREEN),
    ]
    for i, (num, step, clr) in enumerate(steps):
        sy = Inches(1.3) + i * Inches(0.62)
        card(sl, Inches(0.3), sy, Inches(2.5), Inches(0.52))
        # Number circle
        c = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.4), sy + Inches(0.06), Inches(0.38), Inches(0.38))
        c.fill.solid(); c.fill.fore_color.rgb = clr; c.line.fill.background()
        txt(sl, Inches(0.4), sy + Inches(0.06), Inches(0.38), Inches(0.38),
            num, 13, WHITE, bold=True, align=PP_ALIGN.CENTER)
        txt(sl, Inches(0.9), sy + Inches(0.08), Inches(1.8), Inches(0.36), step, 11, LGRAY)

    # Add summary steps on the right too
    highlights = [
        ("12 Categories", "water, roads, electricity...", BLUE),
        ("0-100 Score", "5-factor severity algorithm", ORANGE),
        ("5 Levels", "Ward Officer → Commissioner", RED),
        ("10+ Languages", "Hindi-first, multilingual", GREEN),
    ]
    for i, (title, desc, clr) in enumerate(highlights):
        hy2 = Inches(1.3) + i * Inches(1.45)
        card(sl, Inches(10.5), hy2, Inches(2.5), Inches(1.2))
        txt(sl, Inches(10.7), hy2 + Inches(0.12), Inches(2.1), Inches(0.4), title, 16, clr, bold=True, align=PP_ALIGN.CENTER)
        txt(sl, Inches(10.7), hy2 + Inches(0.55), Inches(2.1), Inches(0.5), desc, 11, MGRAY, align=PP_ALIGN.CENTER)


def slide_06_wireframes(prs):
    """Slide 6: Wireframes & Screen Flow"""
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    slide_header(sl, "Wireframes & Screen Flow")

    # Screen flow diagram (top)
    add_img(sl, "screen_flow", Inches(0.3), Inches(1.2), Inches(12.7), Inches(3.5))

    # Screen descriptions (bottom)
    screens = [
        ("Home Page",        "Hero, live stats counter, 3 CTAs,\nfeature cards, language selector", BLUE),
        ("File Complaint",   "4-step wizard: category grid,\nphoto dropzone, map, review", GREEN),
        ("Voice Assistant",  "Mic button, connection status,\nlive transcript, call timer", PURPLE),
        ("Public Dashboard", "Google Maps heatmap, filter chips,\nstats bar, ward scorecards", BLUE),
        ("Track Complaint",  "Search by phone/#, status card,\ntimeline, legal rights", GREEN),
        ("Admin Queue",      "Filterable table, severity badges,\nSLA countdown, pagination", RED),
        ("Grievance Detail", "Split: complaint + photos |\ntimeline + actions", ORANGE),
    ]
    for i, (name, desc, clr) in enumerate(screens):
        cx = Inches(0.3) + i * Inches(1.85)
        card(sl, cx, Inches(5.0), Inches(1.75), Inches(2.2))
        txt(sl, cx + Inches(0.08), Inches(5.1), Inches(1.6), Inches(0.35),
            name, 12, clr, bold=True, align=PP_ALIGN.CENTER)
        txt(sl, cx + Inches(0.08), Inches(5.5), Inches(1.6), Inches(1.4),
            desc, 10, LGRAY, align=PP_ALIGN.CENTER)


def slide_07_architecture(prs):
    """Slide 7: Architecture Diagram (full slide)"""
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    slide_header(sl, "System Architecture")

    add_img(sl, "architecture", Inches(0.2), Inches(1.15), Inches(12.9), Inches(6.1))


def slide_08_technologies(prs):
    """Slide 8: Technologies Used"""
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    slide_header(sl, "Technologies Used")

    groups = [
        ("Frontend", BLUE, [
            "Next.js 15 (App Router)", "React 19", "Tailwind CSS 4",
            "TypeScript", "LiveKit Client SDK", "Google Maps JS API", "NextAuth.js"
        ]),
        ("Backend", GREEN, [
            "Express.js", "Prisma ORM", "PostgreSQL 14+",
            "PostGIS 3", "Zod Validation", "Node-Cron", "Helmet + CORS"
        ]),
        ("AI / ML", ORANGE, [
            "Gemini 2.0 Flash", "Gemini Vision (Photos)",
            "Gemini 2.5 Flash (Voice)", "5-Factor Severity Algo",
            "Semantic Deduplication", "Legal Rights Generation"
        ]),
        ("Voice Agent", PURPLE, [
            "LiveKit Agents (Python)", "Deepgram STT (Hindi)",
            "ElevenLabs TTS (Hindi)", "Silero VAD",
            "Multilingual Turn Detection"
        ]),
        ("Infrastructure", RED, [
            "Vercel (Frontend)", "Render (API + Agent)",
            "Supabase (DB + Storage)", "LiveKit Cloud (WebRTC)",
            "npm Workspaces (Monorepo)"
        ]),
    ]

    for i, (title, clr, techs) in enumerate(groups):
        cx = Inches(0.35) + i * Inches(2.58)
        card(sl, cx, Inches(1.2), Inches(2.45), Inches(4.8))
        rect(sl, cx, Inches(1.2), Inches(2.45), Inches(0.5), clr)
        txt(sl, cx, Inches(1.22), Inches(2.45), Inches(0.45),
            title, 15, WHITE, bold=True, align=PP_ALIGN.CENTER)
        for j, tech in enumerate(techs):
            txt(sl, cx + Inches(0.15), Inches(1.85) + j * Inches(0.4), Inches(2.15), Inches(0.35),
                f"  {tech}", 11, LGRAY)

    # AWS mapping row
    card(sl, Inches(0.35), Inches(6.15), Inches(12.6), Inches(1.15), DBLUE)
    txt(sl, Inches(0.6), Inches(6.2), Inches(4), Inches(0.3),
        "AWS Services Mapping", 14, ORANGE, bold=True)

    aws = [
        "Gemini  →  Amazon Bedrock",
        "Deepgram  →  Amazon Transcribe",
        "ElevenLabs  →  Amazon Polly",
        "Supabase DB  →  Amazon RDS",
        "Supabase Storage  →  Amazon S3",
        "Vercel  →  AWS Amplify",
        "Render  →  App Runner / ECS",
        "LiveKit  →  Amazon Chime SDK",
    ]
    for i, item in enumerate(aws):
        col, row = i % 4, i // 4
        txt(sl, Inches(0.6) + col * Inches(3.15), Inches(6.55) + row * Inches(0.3),
            Inches(3), Inches(0.28), f"  {item}", 11, LGRAY)


def slide_09_cost(prs):
    """Slide 9: Estimated Implementation Cost"""
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    slide_header(sl, "Estimated Implementation Cost (AWS)")

    # Cost table
    costs = [
        ("Amazon Bedrock (AI)",      "~50K classifications/mo",  "$150 - $300"),
        ("Amazon Transcribe (STT)",  "~500 hrs audio/mo",        "$120"),
        ("Amazon Polly (TTS)",       "~2M characters/mo",        "$8"),
        ("Amazon RDS + PostGIS",     "db.t3.medium, 100GB",      "$70"),
        ("Amazon S3 (Storage)",      "50GB + transfers",         "$5"),
        ("AWS Amplify (Frontend)",   "Standard hosting",         "$15"),
        ("AWS App Runner (API)",     "2 vCPU, 4GB RAM",          "$50"),
        ("Amazon Chime SDK (Voice)", "~10K minutes/mo",          "$30"),
        ("Amazon Location Service",  "~100K geocoding requests", "$50"),
        ("EventBridge + CloudWatch", "Cron + monitoring",        "$11"),
    ]

    headers = ["AWS Service", "Usage Estimate", "Cost / Month"]
    cxs = [Inches(0.5), Inches(4.0), Inches(7.3)]
    cws = [Inches(3.4), Inches(3.2), Inches(1.8)]
    hy = Inches(1.25)

    for hdr, cx, cw in zip(headers, cxs, cws):
        rect(sl, cx, hy, cw, Inches(0.42), BLUE)
        txt(sl, cx + Inches(0.1), hy + Inches(0.05), cw - Inches(0.2), Inches(0.32),
            hdr, 12, WHITE, bold=True, align=PP_ALIGN.CENTER)

    for i, (svc, usage, cost) in enumerate(costs):
        ry = hy + Inches(0.47) + i * Inches(0.42)
        bg = BG_CARD if i % 2 == 0 else BG_CARD2
        for cx, cw, text, clr in zip(cxs, cws, [svc, usage, cost], [WHITE, MGRAY, GREEN]):
            rect(sl, cx, ry, cw, Inches(0.4), bg)
            txt(sl, cx + Inches(0.1), ry + Inches(0.05), cw - Inches(0.2), Inches(0.3),
                text, 11, clr, bold=(clr == GREEN))

    # Total row
    ty = hy + Inches(0.47) + len(costs) * Inches(0.42) + Inches(0.08)
    card(sl, Inches(0.5), ty, Inches(8.6), Inches(0.5), DBLUE)
    txt(sl, Inches(0.7), ty + Inches(0.05), Inches(5), Inches(0.4),
        "TOTAL ESTIMATED (Pilot — 1 City)", 13, WHITE, bold=True)
    txt(sl, Inches(7.3), ty + Inches(0.05), Inches(1.8), Inches(0.4),
        "$509 - $659/mo", 15, GREEN, bold=True, align=PP_ALIGN.CENTER)

    # Scaling cards (right)
    txt(sl, Inches(9.5), Inches(1.25), Inches(3.5), Inches(0.4),
        "Scaling Projections", 18, ORANGE, bold=True)
    accent_line(sl, Inches(9.5), Inches(1.65), Inches(2), ORANGE)

    scales = [
        ("Pilot (1 City)",      "10K users · 5K complaints",   "~$500/mo",    GREEN),
        ("Regional (1 State)",  "100K users · 50K complaints",  "~$2,500/mo",  ORANGE),
        ("National (Pan-India)","1M+ users · 500K+ complaints", "~$15,000/mo", RED),
    ]
    for i, (scale, desc, cost, clr) in enumerate(scales):
        sy = Inches(1.9) + i * Inches(1.65)
        card(sl, Inches(9.5), sy, Inches(3.5), Inches(1.4))
        txt(sl, Inches(9.7), sy + Inches(0.1), Inches(3.1), Inches(0.3), scale, 14, clr, bold=True)
        txt(sl, Inches(9.7), sy + Inches(0.42), Inches(3.1), Inches(0.3), desc, 11, MGRAY)
        txt(sl, Inches(9.7), sy + Inches(0.78), Inches(3.1), Inches(0.4), cost, 20, clr, bold=True)

    # Optimization strategies
    txt(sl, Inches(9.5), Inches(6.85), Inches(3.5), Inches(0.3),
        "Cost Optimization", 13, BLUE, bold=True)
    for i, opt in enumerate(["Bedrock batch inference", "S3 Intelligent Tiering", "RDS Reserved Instances (40% off)", "Lambda for cron (pay-per-use)"]):
        txt(sl, Inches(9.5), Inches(7.1) + i * Inches(0.2), Inches(3.5), Inches(0.2),
            f"  {opt}", 9, MGRAY)


def slide_10_voice_escalation(prs):
    """Slide 10: Voice Agent + Escalation deep dive"""
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    slide_header(sl, "Deep Dive: Voice Agent & Escalation Engine")

    # Left: Voice sequence
    txt(sl, Inches(0.5), Inches(1.1), Inches(6), Inches(0.35),
        "Voice Agent — Sequence Diagram", 16, ORANGE, bold=True)
    add_img(sl, "voice_sequence", Inches(0.2), Inches(1.5), Inches(6.5), Inches(5.7))

    # Right: Escalation flow
    txt(sl, Inches(7.2), Inches(1.1), Inches(5.5), Inches(0.35),
        "Auto-Escalation Engine", 16, RED, bold=True)
    add_img(sl, "escalation_flow", Inches(7), Inches(1.5), Inches(6), Inches(5.7))


def slide_11_database_severity(prs):
    """Slide 11: Database ER + Severity Pie"""
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    slide_header(sl, "Database Schema & Severity Scoring")

    # Left: ER diagram
    txt(sl, Inches(0.5), Inches(1.1), Inches(5), Inches(0.35),
        "Entity Relationship Diagram", 16, GREEN, bold=True)
    add_img(sl, "er_diagram", Inches(0.2), Inches(1.5), Inches(6), Inches(5.7))

    # Right: Severity pie
    txt(sl, Inches(7.2), Inches(1.1), Inches(5.5), Inches(0.35),
        "Severity Score Algorithm (0-100)", 16, ORANGE, bold=True)
    add_img(sl, "severity_pie", Inches(7.2), Inches(1.5), Inches(5.5), Inches(3.5))

    # Severity detail cards below pie
    factors = [
        ("Issue Type Base", "30%", "Water, sewage, roads\nget higher base scores", BLUE),
        ("Affected Population", "25%", "Community issues\nand duplicate count", GREEN),
        ("Vulnerability Index", "20%", "Elderly, disabled,\nBPL, pregnant", ORANGE),
        ("Time Sensitivity", "15%", "Seasonal urgency\nand health risks", RED),
        ("Recurrence", "10%", "Same issue reported\nmultiple times", PURPLE),
    ]
    for i, (name, pct, desc, clr) in enumerate(factors):
        fx = Inches(7.0) + (i % 3) * Inches(2.1)
        fy = Inches(5.3) + (i // 3) * Inches(1.15)
        card(sl, fx, fy, Inches(2), Inches(1.05))
        txt(sl, fx + Inches(0.08), fy + Inches(0.05), Inches(1.3), Inches(0.25), name, 10, clr, bold=True)
        txt(sl, fx + Inches(1.4), fy + Inches(0.05), Inches(0.5), Inches(0.25), pct, 11, WHITE, bold=True, align=PP_ALIGN.RIGHT)
        txt(sl, fx + Inches(0.08), fy + Inches(0.35), Inches(1.8), Inches(0.6), desc, 9, MGRAY)


def slide_12_impact_roadmap(prs):
    """Slide 12: Impact + Deployment + Roadmap"""
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    slide_header(sl, "Impact Metrics & Deployment")

    # Left: Impact metrics
    txt(sl, Inches(0.5), Inches(1.15), Inches(5), Inches(0.35),
        "Projected Impact", 18, ORANGE, bold=True)

    metrics = [
        ("Resolution Time",     "45 days  →  7-14 days",          GREEN),
        ("Citizen Awareness",   "67% unaware  →  90%+ informed",  BLUE),
        ("Misrouted Complaints","~40%  →  <5%",                   ORANGE),
        ("Duplicate Complaints","~30%  →  <10%",                  GREEN),
        ("Dept Accountability", "Near zero  →  100% tracked",     RED),
        ("Low-Literacy Access", "Excluded  →  Fully included",    BLUE),
    ]
    for i, (metric, change, clr) in enumerate(metrics):
        my = Inches(1.6) + i * Inches(0.65)
        card(sl, Inches(0.5), my, Inches(5.7), Inches(0.55))
        txt(sl, Inches(0.7), my + Inches(0.08), Inches(2.3), Inches(0.4), metric, 13, WHITE, bold=True)
        txt(sl, Inches(3.0), my + Inches(0.08), Inches(3), Inches(0.4), change, 13, clr)

    # Government alignment
    txt(sl, Inches(0.5), Inches(5.7), Inches(5), Inches(0.3),
        "Government Alignment", 14, GREEN, bold=True)
    aligns = ["Digital India", "CPGRAMS", "Smart Cities Mission", "RTI Act 2005", "Article 21"]
    for i, a in enumerate(aligns):
        col, row = i % 3, i // 3
        txt(sl, Inches(0.5) + col * Inches(2), Inches(6.05) + row * Inches(0.3),
            Inches(2), Inches(0.28), f"  {a}", 11, LGRAY)

    # Right: Deployment diagram
    txt(sl, Inches(6.8), Inches(1.15), Inches(6), Inches(0.35),
        "Deployment Architecture", 18, BLUE, bold=True)
    add_img(sl, "deployment", Inches(6.5), Inches(1.5), Inches(6.5), Inches(5.5))


def slide_13_closing(prs):
    """Slide 13: Thank You / Closing"""
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)

    accent_line(sl, Inches(0), Inches(0), SLIDE_W, BLUE)

    txt(sl, Inches(1), Inches(1.5), Inches(11.3), Inches(1),
        "JanSunwai AI", 60, BLUE, bold=True, align=PP_ALIGN.CENTER)

    txt(sl, Inches(1.5), Inches(2.8), Inches(10.3), Inches(0.7),
        '"Your Voice Matters. We Make It Heard."', 24, GREEN, align=PP_ALIGN.CENTER)

    accent_line(sl, Inches(5.5), Inches(3.7), Inches(2.3), ORANGE)

    txt(sl, Inches(1.5), Inches(4.1), Inches(10.3), Inches(1),
        "An AI-powered platform transforming how 1.4 billion Indians\n"
        "resolve civic grievances — making government accountable,\n"
        "citizens empowered, and cities smarter.", 17, LGRAY, align=PP_ALIGN.CENTER)

    # Tech badges
    techs = ["Next.js 15", "Express.js", "Gemini AI", "LiveKit", "PostGIS", "AWS-Ready"]
    bw = Inches(1.7)
    gap = Inches(0.2)
    total = len(techs) * bw + (len(techs) - 1) * gap
    sx = (SLIDE_W - total) // 2

    for i, tech in enumerate(techs):
        bx = sx + i * (bw + gap)
        card(sl, bx, Inches(5.5), bw, Inches(0.55), DBLUE)
        txt(sl, bx, Inches(5.54), bw, Inches(0.46), tech, 13, WHITE, align=PP_ALIGN.CENTER)

    txt(sl, Inches(3), Inches(6.5), Inches(7.3), Inches(0.5),
        "Thank You!  |  AWS AI for Bharat Hackathon", 16, MGRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════

def main():
    print("=" * 55)
    print("  Building JanSunwai AI PPTX")
    print("=" * 55)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        ("01 Title",                    slide_01_title),
        ("02 Brief About Idea",         slide_02_idea),
        ("03 Differentiation & USP",    slide_03_usp),
        ("04 Features + Mindmap",       slide_04_features),
        ("05 Process Flow",             slide_05_process_flow),
        ("06 Wireframes / Screen Flow", slide_06_wireframes),
        ("07 Architecture",             slide_07_architecture),
        ("08 Technologies",             slide_08_technologies),
        ("09 Cost Estimation",          slide_09_cost),
        ("10 Voice Agent & Escalation", slide_10_voice_escalation),
        ("11 Database & Severity",      slide_11_database_severity),
        ("12 Impact & Deployment",      slide_12_impact_roadmap),
        ("13 Thank You",                slide_13_closing),
    ]

    for name, builder in builders:
        print(f"  Building {name}...")
        builder(prs)

    prs.save(OUTPUT_FILE)
    print(f"\n  Saved: {OUTPUT_FILE}")
    print(f"  Slides: {len(prs.slides)}")
    print("=" * 55)


if __name__ == "__main__":
    main()
